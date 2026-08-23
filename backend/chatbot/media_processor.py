import os
import csv
import zipfile
import xml.etree.ElementTree as ET
from PIL import Image
import pytesseract
from striprtf.striprtf import rtf_to_text
import docx
from pptx import Presentation
import openpyxl

from config.environment import MEDIA_PATH

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_csv(file_path: str) -> str:
    text_lines = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        headers = next(reader, None)
        for i, row in enumerate(reader, start=1):
            row_str = f"Row {i}: "
            if headers:
                row_str += ", ".join(f"{h}: {val}" for h, val in zip(headers, row))
            else:
                row_str += ", ".join(row)
            text_lines.append(row_str)
    return "\n".join(text_lines)

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"Error parsing DOCX {file_path}: {e}")
        # Fallback raw extraction
        try:
            with zipfile.ZipFile(file_path) as docx_zip:
                xml_content = docx_zip.read('word/document.xml')
                root = ET.fromstring(xml_content)
                namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                paragraphs = []
                for p in root.findall('.//w:p', namespaces):
                    texts = [t.text for t in p.findall('.//w:t', namespaces) if t.text]
                    if texts:
                        paragraphs.append("".join(texts))
                return "\n".join(paragraphs)
        except Exception as fallback_err:
            print(f"DOCX fallback parse failed: {fallback_err}")
            return ""

def extract_text_from_xlsx(file_path: str) -> str:
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets_text = []
        for name in wb.sheetnames:
            ws = wb[name]
            sheet_text = [f"--- Sheet: {name} ---"]
            for r_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                row_vals = [str(val) if val is not None else "" for val in row]
                if any(row_vals):
                    sheet_text.append(f"Row {r_idx}: " + ", ".join(row_vals))
            sheets_text.append("\n".join(sheet_text))
        return "\n\n".join(sheets_text)
    except Exception as e:
        print(f"Error parsing XLSX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides_text.append("\n".join(slide_text))
        return "\n\n".join(slides_text)
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")
        return ""

def extract_text_from_rtf(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)
    except Exception as e:
        print(f"Error parsing RTF {file_path}: {e}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error parsing Image OCR {file_path}: {e}")
        return ""

def extract_text_from_audio(file_path: str) -> str:
    try:
        import speech_recognition as sr
        r = sr.Recognizer()
        
        ext = os.path.splitext(file_path)[1].lower()
        target_wav = file_path
        temp_wav = None
        
        # If it is not a WAV file, try converting it to WAV via pydub
        if ext != ".wav":
            try:
                from pydub import AudioSegment
                sound = AudioSegment.from_file(file_path)
                temp_wav = file_path + ".temp.wav"
                sound.export(temp_wav, format="wav")
                target_wav = temp_wav
            except Exception as conv_err:
                print(f"Audio conversion failed (requires pydub + ffmpeg): {conv_err}")
                return f"[Audio File: {os.path.basename(file_path)} - Conversion to WAV failed: {conv_err}]"
        
        with sr.AudioFile(target_wav) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
            
        if temp_wav and os.path.exists(temp_wav):
            try:
                os.remove(temp_wav)
            except Exception:
                pass
                
        return f"[Audio Transcription of {os.path.basename(file_path)}]\n{text}"
    except Exception as e:
        print(f"Audio transcription error {file_path}: {e}")
        return f"[Audio File: {os.path.basename(file_path)} - Transcription failed: {e}]"

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        from chatbot.pdf_indexer import extract_text_from_pdf
        return extract_text_from_pdf(file_path)
    elif ext in [".txt", ".log", ".json", ".xml", ".html"]:
        return extract_text_from_txt(file_path)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    elif ext == ".docx" or ext == ".doc":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx" or ext == ".xls":
        return extract_text_from_xlsx(file_path)
    elif ext == ".pptx" or ext == ".ppt":
        return extract_text_from_pptx(file_path)
    elif ext == ".rtf":
        return extract_text_from_rtf(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
        return extract_text_from_image(file_path)
    elif ext in [".wav", ".mp3", ".m4a", ".mp4"]:
        return extract_text_from_audio(file_path)
    else:
        # Fallback to plain text read
        try:
            return extract_text_from_txt(file_path)
        except Exception:
            raise ValueError(f"Unsupported file format: {ext}")
